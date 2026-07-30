# -*- coding: utf-8 -*-
"""
报告生成与可视化导出模块
支持将爬取的报告/研报/公告整理分析后导出为 PPT、PDF、Word、Excel 格式
"""

import os
import re
import json
import time
from pathlib import Path

# 标准库 fallback 生成器（当 python-pptx/docx/openpyxl 未安装时使用）
try:
    from .report_stdlib_fallbacks import StdlibPptxWriter, StdlibDocxWriter, StdlibXlsxWriter
except ImportError:
    from report_stdlib_fallbacks import StdlibPptxWriter, StdlibDocxWriter, StdlibXlsxWriter
from typing import Dict, Any, List, Optional, Union
from datetime import datetime
from dataclasses import dataclass, field

SKILL_DATA_DIR = Path(__file__).parent.parent / "data"
EXPORT_DIR = SKILL_DATA_DIR / "exports"


@dataclass
class ExportedReport:
    """导出报告信息"""
    title: str
    file_path: str
    file_type: str  # ppt, pdf, docx, xlsx
    created_at: str = ""
    stock_code: str = ""
    stock_name: str = ""


class DataOrganizer:
    """数据整理器 - 将原始爬取数据整理成结构化报告"""

    @staticmethod
    def organize_periodic_reports(reports: List) -> Dict[str, Any]:
        """整理定期报告"""
        result = {
            "annual": [],      # 年报
            "half_year": [],   # 半年报
            "quarter": [],     # 季报
            "other": []        # 其他
        }

        for r in reports:
            title = getattr(r, 'title', '') or ''
            ptype = getattr(r, 'report_type', '') or ''

            if '年报' in ptype or 'annual' in ptype.lower():
                result["annual"].append({
                    "title": title,
                    "date": getattr(r, 'publish_date', ''),
                    "url": getattr(r, 'url', ''),
                    "stock_code": getattr(r, 'stock_code', '')
                })
            elif '半年' in ptype or 'half' in ptype.lower():
                result["half_year"].append({
                    "title": title,
                    "date": getattr(r, 'publish_date', ''),
                    "url": getattr(r, 'url', ''),
                    "stock_code": getattr(r, 'stock_code', '')
                })
            elif '季' in ptype or 'quarter' in ptype.lower():
                result["quarter"].append({
                    "title": title,
                    "date": getattr(r, 'publish_date', ''),
                    "url": getattr(r, 'url', '')
                })
            else:
                result["other"].append({
                    "title": title,
                    "date": getattr(r, 'publish_date', ''),
                    "url": getattr(r, 'url', '')
                })

        return result

    @staticmethod
    def organize_broker_reports(reports: List) -> Dict[str, Any]:
        """整理券商研报"""
        result = {
            "buy": [],      # 买入评级
            "increase": [],  # 增持
            "neutral": [],   # 中性
            "reduce": [],    # 减持
            "statistics": {
                "total": len(reports),
                "brokers": {},
                "analysts": {},
                "rating_distribution": {}
            }
        }

        for r in reports:
            rating = getattr(r, 'rating', '') or '未知'
            broker = getattr(r, 'broker_name', '') or '未知券商'
            analyst = getattr(r, 'analyst', '') or '未知分析师'

            # 统计
            result["statistics"]["brokers"][broker] = result["statistics"]["brokers"].get(broker, 0) + 1
            result["statistics"]["analysts"][analyst] = result["statistics"]["analysts"].get(analyst, 0) + 1

            # 按评级分类
            report_dict = {
                "title": getattr(r, 'title', ''),
                "date": getattr(r, 'publish_date', ''),
                "broker": broker,
                "analyst": analyst,
                "rating": rating,
                "target_price": getattr(r, 'target_price', 0),
                "target_change_pct": getattr(r, 'target_change_pct', 0),
                "url": getattr(r, 'url', '')
            }

            if '买入' in rating or '推荐' in rating:
                result["buy"].append(report_dict)
            elif '增持' in rating:
                result["increase"].append(report_dict)
            elif '中性' in rating:
                result["neutral"].append(report_dict)
            else:
                result["reduce"].append(report_dict)

        # 评级分布
        result["statistics"]["rating_distribution"] = {
            "买入": len(result["buy"]),
            "增持": len(result["increase"]),
            "中性": len(result["neutral"]),
            "减持": len(result["reduce"])
        }

        return result

    @staticmethod
    def organize_announcements(announcements: List) -> Dict[str, Any]:
        """整理公告"""
        result = {
            "by_type": {},
            "recent": [],
            "statistics": {
                "total": len(announcements)
            }
        }

        for a in announcements:
            title = getattr(a, 'title', '') or ''
            ann_type = getattr(a, 'announcement_type', '') or '其他'
            date = getattr(a, 'publish_date', '') or ''

            if ann_type not in result["by_type"]:
                result["by_type"][ann_type] = []
            result["by_type"][ann_type].append({
                "title": title,
                "date": date,
                "url": getattr(a, 'url', '')
            })

            # 最近公告
            if len(result["recent"]) < 20:
                result["recent"].append({
                    "title": title,
                    "date": date,
                    "type": ann_type,
                    "url": getattr(a, 'url', '')
                })

        return result

    @staticmethod
    def generate_summary(data: Dict[str, Any], stock_code: str = "") -> Dict[str, Any]:
        """生成综合摘要"""
        summary = {
            "stock_code": stock_code,
            "generated_at": datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            "periodic": {"count": 0, "latest": ""},
            "broker": {"count": 0, "buy_count": 0, "latest": ""},
            "announcement": {"count": 0, "latest": ""}
        }

        # 定期报告摘要
        if "periodic_reports" in data:
            reports = data["periodic_reports"]
            summary["periodic"]["count"] = len(reports)
            dates = [getattr(r, 'publish_date', '') for r in reports if hasattr(r, 'publish_date')]
            summary["periodic"]["latest"] = max(dates) if dates else ""

        # 券商研报摘要
        if "broker_reports" in data:
            reports = data["broker_reports"]
            summary["broker"]["count"] = len(reports)
            buy_count = sum(1 for r in reports if getattr(r, 'rating', '') in ['买入', '推荐', '增持'])
            summary["broker"]["buy_count"] = buy_count
            dates = [getattr(r, 'publish_date', '') for r in reports if hasattr(r, 'publish_date')]
            summary["broker"]["latest"] = max(dates) if dates else ""

        # 公告摘要
        if "announcements" in data:
            announcements = data["announcements"]
            summary["announcement"]["count"] = len(announcements)
            dates = [getattr(a, 'publish_date', '') for a in announcements if hasattr(a, 'publish_date')]
            summary["announcement"]["latest"] = max(dates) if dates else ""

        return summary


class AnalysisEngine:
    """分析引擎 - 对整理后的数据进行分析"""

    @staticmethod
    def analyze_broker_rating_trend(reports: List) -> Dict[str, Any]:
        """分析券商评级趋势"""
        if not reports:
            return {}

        # 按时间排序
        sorted_reports = sorted(reports, key=lambda x: getattr(x, 'publish_date', '') or '')

        rating_map = {"买入": 5, "增持": 4, "中性": 3, "减持": 2, "卖出": 1}

        trend = {
            "average_rating": 0,
            "rating_distribution": {},
            "target_price_stats": {},
            "latest_rating": ""
        }

        ratings = []
        target_prices = []
        target_changes = []

        for r in sorted_reports:
            rating = getattr(r, 'rating', '') or '未知'
            rating_val = rating_map.get(rating, 3)
            ratings.append(rating_val)

            target_price = getattr(r, 'target_price', 0)
            if target_price > 0:
                target_prices.append(target_price)

            target_change = getattr(r, 'target_change_pct', 0)
            if target_change != 0:
                target_changes.append(target_change)

            trend["rating_distribution"][rating] = trend["rating_distribution"].get(rating, 0) + 1

        if ratings:
            trend["average_rating"] = sum(ratings) / len(ratings)
            trend["latest_rating"] = sorted_reports[-1].rating if hasattr(sorted_reports[-1], 'rating') else ""

        if target_prices:
            trend["target_price_stats"] = {
                "average": sum(target_prices) / len(target_prices),
                "max": max(target_prices),
                "min": min(target_prices),
                "count": len(target_prices)
            }

        if target_changes:
            trend["average_target_change"] = sum(target_changes) / len(target_changes)

        return trend

    @staticmethod
    def analyze_announcement_pattern(announcements: List) -> Dict[str, Any]:
        """分析公告模式"""
        if not announcements:
            return {}

        pattern = {
            "total_count": len(announcements),
            "by_month": {},
            "by_type": {},
            "high_frequency_types": []
        }

        for a in announcements:
            date = getattr(a, 'publish_date', '') or ''
            if len(date) >= 7:
                month = date[:7]  # YYYY-MM
                pattern["by_month"][month] = pattern["by_month"].get(month, 0) + 1

            ann_type = getattr(a, 'announcement_type', '') or '其他'
            pattern["by_type"][ann_type] = pattern["by_type"].get(ann_type, 0) + 1

        # 高频类型排序
        pattern["high_frequency_types"] = sorted(
            pattern["by_type"].items(),
            key=lambda x: -x[1]
        )[:5]

        return pattern


class PPTExporter:
    """PPT导出器"""

    # v4.3.2 修复：universal 中文字体
    DEFAULT_FONT_LATIN = "Times New Roman"
    DEFAULT_FONT_EA = "宋体"

    def __init__(self):
        self.pptx_available = False
        self._fallback = StdlibPptxWriter()
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RgbColor
            self.pptx_available = True
            self.Presentation = Presentation
            self.Inches = Inches
            self.Pt = Pt
            self.Emu = Emu
            self.RgbColor = RgbColor
        except ImportError:
            print("[警告] python-pptx 未安装，使用标准库 fallback 生成 PPT")

    def _set_font(self, run, size: int = 18, bold: bool = False) -> None:
        """v4.3.2: 设置文本 run 的中文字体。"""
        if not self.pptx_available:
            return
        run.font.name = self.DEFAULT_FONT_LATIN
        run.font.size = self.Pt(size)
        run.font.bold = bold
        # 设置东亚字体（通过 XML）
        rPr = run._r.get_or_add_rPr()
        ea = rPr.makeelement(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
            {"typeface": self.DEFAULT_FONT_EA}
        )
        latin = rPr.makeelement(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}latin",
            {"typeface": self.DEFAULT_FONT_LATIN}
        )
        # 替换已有的 ea/latin 元素
        for old in rPr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}ea"):
            rPr.remove(old)
        for old in rPr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}latin"):
            rPr.remove(old)
        rPr.insert(0, ea)
        rPr.insert(0, latin)

    def create_presentation(self) -> Any:
        """创建演示文稿"""
        if not self.pptx_available:
            return None
        return self.Presentation()

    def add_title_slide(self, prs, title: str, subtitle: str = ""):
        """添加标题幻灯片"""
        if not self.pptx_available:
            return

        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = ""
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        self._set_font(run, size=36, bold=True)

        if subtitle and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = ""
            tf = subtitle_shape.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = subtitle
            self._set_font(run, size=18, bold=False)

    def add_content_slide(self, prs, title: str, content: List[str], columns: int = 1):
        """添加内容幻灯片"""
        if not self.pptx_available:
            return

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # 标题
        title_shape = slide.shapes.title
        title_shape.text = ""
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        self._set_font(run, size=28, bold=True)

        # 添加内容
        left = self.Inches(0.5)
        top = self.Inches(1.5)
        width = self.Inches(9)
        height = self.Inches(0.5)

        for i, text in enumerate(content):
            textbox = slide.shapes.add_textbox(left, top + self.Inches(i * 0.5), width, height)
            textframe = textbox.text_frame
            textframe.word_wrap = True  # v4.3.2: 修复长文本排版
            p = textframe.paragraphs[0]
            run = p.add_run()
            run.text = text
            self._set_font(run, size=14, bold=False)

    def export_comprehensive_report(self, data: Dict[str, Any],
                                     stock_code: str,
                                     output_path: str = "") -> str:
        """
        导出综合报告为PPT

        Args:
            data: 综合报告数据
            stock_code: 股票代码
            output_path: 输出路径

        Returns:
            保存路径
        """
        if not self.pptx_available:
            return self._fallback.export_comprehensive_report(data, stock_code, output_path)

        prs = self.create_presentation()

        summary = data.get("summary", {})
        stock_name = summary.get("stock_name", stock_code)

        # 标题页
        self.add_title_slide(
            prs,
            f"{stock_name} ({stock_code}) 综合报告",
            f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        )

        # 概览页
        overview = [
            f"定期报告: {summary.get('periodic_count', 0)} 份",
            f"券商研报: {summary.get('broker_count', 0)} 份",
            f"公告: {summary.get('announcement_count', 0)} 份",
            f"买入评级研报: {'是' if summary.get('has_buy_rating') else '否'}",
            f"最新定期报告: {summary.get('latest_periodic_date', 'N/A')}",
            f"最新券商研报: {summary.get('latest_broker_date', 'N/A')}",
            f"最新公告: {summary.get('latest_announcement_date', 'N/A')}"
        ]
        self.add_content_slide(prs, "报告概览", overview)

        # 定期报告页
        periodic = data.get("periodic_reports", [])
        if periodic:
            periodic_content = [
                f"共 {len(periodic)} 份定期报告",
                ""
            ]
            for p in periodic[:10]:
                title = getattr(p, 'title', '')[:50]
                date = getattr(p, 'publish_date', '')[:10]
                url = getattr(p, 'url', '')
                line = f"• {date} - {title}"
                if url:
                    line += f"\n   {url}"  # v4.4.0: PPT 添加原文链接
                periodic_content.append(line)
            self.add_content_slide(prs, "定期报告", periodic_content)

        # 券商研报页
        broker = data.get("broker_reports", [])
        if broker:
            # 统计
            rating_stats = {"买入": 0, "增持": 0, "中性": 0, "减持": 0}
            for b in broker:
                rating = getattr(b, 'rating', '')
                if rating in rating_stats:
                    rating_stats[rating] += 1

            broker_content = [
                f"共 {len(broker)} 份研报",
                f"评级分布: 买入{rating_stats['买入']} 增持{rating_stats['增持']} 中性{rating_stats['中性']} 减持{rating_stats['减持']}",
                ""
            ]
            for b in broker[:8]:
                title = getattr(b, 'title', '')[:40]
                broker_name = getattr(b, 'broker_name', '')
                rating = getattr(b, 'rating', '')
                url = getattr(b, 'url', '')
                line = f"• [{rating}] {broker_name}: {title}"
                if url:
                    line += f"\n   {url}"  # v4.4.0: PPT 添加原文链接
                broker_content.append(line)
            self.add_content_slide(prs, "券商研报", broker_content)

        # 公告页
        announcements = data.get("announcements", [])
        if announcements:
            ann_content = [
                f"共 {len(announcements)} 条公告",
                ""
            ]
            for a in announcements[:10]:
                title = getattr(a, 'title', '')[:50]
                date = getattr(a, 'publish_date', '')[:10]
                url = getattr(a, 'url', '')
                line = f"• {date} - {title}"
                if url:
                    line += f"\n   {url}"  # v4.4.0: PPT 添加原文链接
                ann_content.append(line)
            self.add_content_slide(prs, "公告", ann_content)

        # 保存
        if not output_path:
            output_path = str(EXPORT_DIR / f"{stock_code}_report.pptx")

        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)

        return output_path


class WordExporter:
    """Word导出器"""

    # v4.3.2 修复：universal 中文字体默认值（Times New Roman + 宋体）
    DEFAULT_FONT_NAME = "宋体"
    DEFAULT_FONT_LATIN = "Times New Roman"
    DEFAULT_FONT_SIZE = 11  # Pt, 在 _setup_default_font 中转换

    def __init__(self):
        self.docx_available = False
        self._fallback = StdlibDocxWriter()
        try:
            from docx import Document
            from docx.shared import Inches, Pt, RGBColor
            self.docx_available = True
            self.Document = Document
            self.Inches = Inches
            self.Pt = Pt
            self.RGBColor = RGBColor
        except ImportError:
            print("[警告] python-docx 未安装，使用标准库 fallback 生成 DOCX")

    def _setup_default_font(self, doc) -> None:
        """v4.3.2: 将默认样式字体设为 宋体（中文）+ Times New Roman（拉丁），
        确保跨平台中文不乱码。"""
        if not self.docx_available:
            return
        style = doc.styles["Normal"]
        font = style.font
        font.name = self.DEFAULT_FONT_LATIN
        font.size = self.Pt(self.DEFAULT_FONT_SIZE)
        # east-Asia 字体通过 rFonts 元素设置
        rPr = style.element.get_or_add_rPr()
        rFonts = rPr.makeelement(
            "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}rFonts",
            {}
        )
        rFonts.set(
            "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}ascii",
            self.DEFAULT_FONT_LATIN
        )
        rFonts.set(
            "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}hAnsi",
            self.DEFAULT_FONT_LATIN
        )
        rFonts.set(
            "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}eastAsia",
            self.DEFAULT_FONT_NAME
        )
        rFonts.set(
            "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}cs",
            self.DEFAULT_FONT_LATIN
        )
        rPr.insert(0, rFonts)

    def create_document(self) -> Any:
        """创建Word文档"""
        if not self.docx_available:
            return None
        doc = self.Document()
        self._setup_default_font(doc)  # v4.3.2: 设置中文字体
        return doc

    def export_comprehensive_report(self, data: Dict[str, Any],
                                     stock_code: str,
                                     output_path: str = "") -> str:
        """
        导出综合报告为Word

        Args:
            data: 综合报告数据
            stock_code: 股票代码
            output_path: 输出路径

        Returns:
            保存路径
        """
        if not self.docx_available:
            return self._fallback.export_comprehensive_report(data, stock_code, output_path)

        doc = self.create_document()

        summary = data.get("summary", {})
        stock_name = summary.get("stock_name", stock_code)

        # 标题
        doc.add_heading(f"{stock_name} ({stock_code}) 综合分析报告", 0)

        # 基本信息
        doc.add_heading("报告概览", 1)
        info_table = doc.add_table(rows=7, cols=2)
        info_table.style = 'Light Grid Accent 1'

        info_data = [
            ("股票代码", stock_code),
            ("股票名称", stock_name),
            ("定期报告", f"{summary.get('periodic_count', 0)} 份"),
            ("券商研报", f"{summary.get('broker_count', 0)} 份"),
            ("公告", f"{summary.get('announcement_count', 0)} 份"),
            ("买入评级研报", "是" if summary.get('has_buy_rating') else "否"),
            ("生成时间", datetime.now().strftime('%Y-%m-%d %H:%M:%S'))
        ]

        for i, (key, val) in enumerate(info_data):
            info_table.rows[i].cells[0].text = key
            info_table.rows[i].cells[1].text = val

        # 定期报告
        periodic = data.get("periodic_reports", [])
        if periodic:
            doc.add_heading(f"定期报告 ({len(periodic)} 份)", 1)
            for p in periodic:
                title = getattr(p, 'title', '')
                date = getattr(p, 'publish_date', '')[:10]
                url = getattr(p, 'url', '')
                doc.add_paragraph(f"• {date} - {title}", style='List Bullet')
                if url:
                    doc.add_paragraph(f"  链接: {url}", style='Normal')

        # 券商研报
        broker = data.get("broker_reports", [])
        if broker:
            doc.add_heading(f"券商研报 ({len(broker)} 份)", 1)

            # 评级统计
            rating_stats = {"买入": 0, "增持": 0, "中性": 0, "减持": 0}
            for b in broker:
                rating = getattr(b, 'rating', '')
                if rating in rating_stats:
                    rating_stats[rating] += 1

            doc.add_paragraph(f"评级分布: 买入{rating_stats['买入']} 增持{rating_stats['增持']} 中性{rating_stats['中性']} 减持{rating_stats['减持']}")

            for b in broker:
                title = getattr(b, 'title', '')
                date = getattr(b, 'publish_date', '')[:10]
                broker_name = getattr(b, 'broker_name', '')
                analyst = getattr(b, 'analyst', '')
                rating = getattr(b, 'rating', '')
                target_price = getattr(b, 'target_price', 0)

                p = doc.add_paragraph()
                p.add_run(f"• {date} [{rating}] {broker_name}").bold = True
                p.add_run(f"\n  {title}")
                if analyst:
                    p.add_run(f"\n  分析师: {analyst}")
                if target_price > 0:
                    p.add_run(f"\n  目标价: ¥{target_price:.2f}")
                # v4.4.0: 添加原文链接
                url = getattr(b, 'url', '')
                if url:
                    p.add_run(f"\n  链接: {url}")

        # 公告
        announcements = data.get("announcements", [])
        if announcements:
            doc.add_heading(f"公告 ({len(announcements)} 条)", 1)
            for a in announcements[:30]:
                title = getattr(a, 'title', '')
                date = getattr(a, 'publish_date', '')[:10]
                ann_type = getattr(a, 'announcement_type', '')
                url = getattr(a, 'url', '')

                p = doc.add_paragraph()
                p.add_run(f"• {date} [{ann_type}] ").bold = True
                p.add_run(title)
                if url:
                    p.add_run(f"\n  链接: {url}")

        # 保存
        if not output_path:
            output_path = str(EXPORT_DIR / f"{stock_code}_report.docx")

        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        doc.save(output_path)

        return output_path

    def export_with_images(self, data: Dict[str, Any],
                           output_path: str = "",
                           max_image_width_inches: float = 5.5) -> str:
        """v4.5 新增：导出带图片的 Word 文档。

        Args:
            data: 文档结构
                {"title": "...",
                 "metadata": {"source": "...", "date": "..."},
                 "sections": [
                     {"heading": "一、xxx", "level": 1,
                      "paragraphs": ["段落1", "段落2"],
                      "images": [{"path": "./a.png", "caption": "图1", "width": 5.5}]
                     },
                     ...
                 ]}
            output_path: 输出路径（默认 data/exports/）
            max_image_width_inches: 图片最大宽度

        Returns:
            保存路径（fallback 时返回 .md/.txt）
        """
        if not self.docx_available:
            # Fallback：直接写 markdown（用 stdlib）
            md_path = (output_path or str(EXPORT_DIR / "export_with_images.md"))
            if not md_path.endswith(".md"):
                md_path = md_path.rsplit(".", 1)[0] + ".md"
            lines = [f"# {data.get('title', '无标题')}\n\n"]
            meta = data.get("metadata", {})
            if meta:
                lines.append(" · ".join(f"**{k}**: {v}" for k, v in meta.items()) + "\n\n")
            for sec in data.get("sections", []):
                lines.append(f"\n## {sec.get('heading', '')}\n\n")
                for p in sec.get("paragraphs", []):
                    lines.append(f"{p}\n\n")
                for img in sec.get("images", []):
                    if isinstance(img, dict):
                        cap = img.get("caption", "")
                        path = img.get("path", "")
                        lines.append(f"![{cap}]({path})\n\n")
                    else:
                        lines.append(f"![]({img})\n\n")
            Path(md_path).parent.mkdir(parents=True, exist_ok=True)
            Path(md_path).write_text("".join(lines), encoding="utf-8")
            return md_path

        doc = self.create_document()
        title = data.get("title", "无标题")
        doc.add_heading(title, 0)

        # 元信息
        meta = data.get("metadata", {})
        if meta:
            doc.add_paragraph(" · ".join(
                f"{k}: {v}" for k, v in meta.items()
            ))

        # 各章节
        for sec in data.get("sections", []):
            level = sec.get("level", 1)
            heading = sec.get("heading", "")
            doc.add_heading(heading, level=min(4, level + 1))

            for para in sec.get("paragraphs", []):
                doc.add_paragraph(para)

            for img in sec.get("images", []):
                if isinstance(img, dict):
                    img_path = img.get("path", "")
                    caption = img.get("caption", "")
                    width = img.get("width", max_image_width_inches)
                else:
                    img_path = img
                    caption = ""
                    width = max_image_width_inches

                if not img_path or not Path(img_path).exists():
                    doc.add_paragraph(f"[图片缺失: {img_path}]")
                    continue
                try:
                    p = doc.add_paragraph()
                    p.alignment = 1  # 居中
                    run = p.add_run()
                    run.add_picture(img_path, width=self.Inches(width))
                    if caption:
                        cap_p = doc.add_paragraph()
                        cap_p.alignment = 1
                        cap_p.add_run(caption).italic = True
                except Exception as e:
                    doc.add_paragraph(f"[图片加载失败: {img_path} — {e}]")

        # 保存
        if not output_path:
            ts = time.strftime("%Y%m%d_%H%M%S")
            output_path = str(EXPORT_DIR / f"{ts}_{title[:30]}_with_images.docx")

        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        doc.save(output_path)
        return output_path


# v4.5 新增：Markdown 导出器（在 report_exporter 里复用 html2md.convert）
class MarkdownExporter:
    """v4.5 新增：Markdown 导出器（与 scripts/markdown_exporter.py 行为一致）"""

    def __init__(self, output_dir: Optional[str] = None):
        self.output_dir = Path(output_dir) if output_dir else EXPORT_DIR
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def export(self, data: Dict[str, Any], output_path: Optional[str] = None,
               images_inline: bool = True) -> str:
        """导出单篇 Markdown。"""
        # 委托给 scripts.markdown_exporter（避免代码重复）
        try:
            from .markdown_exporter import export_to_markdown as _export
            if output_path is None:
                output_path = str(self.output_dir / f"{time.strftime('%Y%m%d_%H%M%S')}.md")
            return _export(data, output_path)
        except ImportError:
            # 纯 stdlib fallback
            if output_path is None:
                output_path = str(self.output_dir / "export.md")
            content = data.get("content", "")
            if isinstance(content, str) and content.lstrip().startswith("<"):
                try:
                    from .html2md import convert
                    content = convert(content)
                except ImportError:
                    content = re.sub(r"<[^>]+>", "", content)
            lines = [f"# {data.get('title', '无标题')}\n\n"]
            if data.get("source"):
                lines.append(f"**来源**: {data['source']}\n\n")
            lines.append(content)
            Path(output_path).parent.mkdir(parents=True, exist_ok=True)
            Path(output_path).write_text("".join(lines), encoding="utf-8")
            return output_path

    def batch(self, items: List[Dict[str, Any]],
              batch_id: Optional[str] = None) -> Dict[str, str]:
        """批量导出。"""
        try:
            from .markdown_exporter import batch_to_markdown as _batch
            return _batch(items, batch_id=batch_id,
                          output_dir=str(self.output_dir))
        except ImportError:
            # fallback
            return {"count": "0"}


class ExcelExporter:
    """Excel导出器"""

    # v4.3.2 修复：universal 中文字体默认值
    DEFAULT_FONT_NAME = "宋体"

    def __init__(self):
        self.openpyxl_available = False
        self._fallback = StdlibXlsxWriter()
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, Alignment, PatternFill
            self.openpyxl_available = True
            self.Workbook = Workbook
            self.Font = Font
            self.Alignment = Alignment
            self.PatternFill = PatternFill
        except ImportError:
            print("[警告] openpyxl 未安装，使用标准库 fallback 生成 XLSX")

    def _apply_font_to_cells(self, ws, start_row: int = 1, end_row: int = 100,
                             start_col: int = 1, end_col: int = 10) -> None:
        """v4.3.2: 为指定区域的所有单元格设置宋体字体。"""
        if not self.openpyxl_available:
            return
        for row in ws.iter_rows(min_row=start_row, max_row=end_row,
                                min_col=start_col, max_col=end_col):
            for cell in row:
                if cell.value is not None:
                    cell.font = self.Font(name=self.DEFAULT_FONT_NAME, size=11)

    def export_comprehensive_report(self, data: Dict[str, Any],
                                     stock_code: str,
                                     output_path: str = "") -> str:
        """
        导出综合报告为Excel

        Args:
            data: 综合报告数据
            stock_code: 股票代码
            output_path: 输出路径

        Returns:
            保存路径
        """
        if not self.openpyxl_available:
            return self._fallback.export_comprehensive_report(data, stock_code, output_path)

        wb = self.Workbook()

        summary = data.get("summary", {})
        stock_name = summary.get("stock_name", stock_code)
        HEADER_FONT = self.Font(size=11, bold=True, name=self.DEFAULT_FONT_NAME)

        # Sheet 1: 概览
        ws_overview = wb.active
        ws_overview.title = "报告概览"

        ws_overview['A1'] = f"{stock_name} ({stock_code}) 综合报告"
        ws_overview['A1'].font = self.Font(size=16, bold=True, name=self.DEFAULT_FONT_NAME)

        overview_data = [
            ("股票代码", stock_code),
            ("股票名称", stock_name),
            ("定期报告数", summary.get('periodic_count', 0)),
            ("券商研报数", summary.get('broker_count', 0)),
            ("公告数", summary.get('announcement_count', 0)),
            ("有买入评级", "是" if summary.get('has_buy_rating') else "否"),
            ("最新定期报告", summary.get('latest_periodic_date', '')),
            ("最新券商研报", summary.get('latest_broker_date', '')),
            ("最新公告", summary.get('latest_announcement_date', '')),
            ("生成时间", datetime.now().strftime('%Y-%m-%d %H:%M:%S'))
        ]

        for i, (key, val) in enumerate(overview_data, 2):
            cell_k = ws_overview.cell(row=i, column=1, value=key)
            cell_v = ws_overview.cell(row=i, column=2, value=val)
            cell_k.font = HEADER_FONT
            cell_v.font = self.Font(name=self.DEFAULT_FONT_NAME, size=11)

        # v4.3.2: 设置概览sheet列宽（排版修复）
        ws_overview.column_dimensions['A'].width = 18
        ws_overview.column_dimensions['B'].width = 60

        # Sheet 2: 定期报告
        periodic = data.get("periodic_reports", [])
        ws_periodic = wb.create_sheet("定期报告")
        headers = ["日期", "标题", "报告类型", "URL"]
        for col, h in enumerate(headers, 1):
            cell = ws_periodic.cell(row=1, column=col, value=h)
            cell.font = HEADER_FONT

        for i, p in enumerate(periodic, 2):
            ws_periodic.cell(row=i, column=1, value=getattr(p, 'publish_date', '')[:10])
            ws_periodic.cell(row=i, column=2, value=getattr(p, 'title', ''))
            ws_periodic.cell(row=i, column=3, value=getattr(p, 'report_type', ''))
            ws_periodic.cell(row=i, column=4, value=getattr(p, 'url', ''))
        # v4.3.2: 设置列宽（排版修复）
        if periodic:
            widths = [14, 50, 16, 50]
            for i, w in enumerate(widths):
                col_letter = chr(65 + i)
                ws_periodic.column_dimensions[col_letter].width = w

        # Sheet 3: 券商研报
        broker = data.get("broker_reports", [])
        if broker:
            ws_broker = wb.create_sheet("券商研报")
            headers = ["日期", "标题", "券商", "分析师", "评级", "目标价", "目标涨幅", "URL"]
            for col, h in enumerate(headers, 1):
                cell = ws_broker.cell(row=1, column=col, value=h)
                cell.font = HEADER_FONT

            for i, b in enumerate(broker, 2):
                ws_broker.cell(row=i, column=1, value=getattr(b, 'publish_date', '')[:10])
                ws_broker.cell(row=i, column=2, value=getattr(b, 'title', ''))
                ws_broker.cell(row=i, column=3, value=getattr(b, 'broker_name', ''))
                ws_broker.cell(row=i, column=4, value=getattr(b, 'analyst', ''))
                ws_broker.cell(row=i, column=5, value=getattr(b, 'rating', ''))
                ws_broker.cell(row=i, column=6, value=float(getattr(b, 'target_price', 0) or 0))
                ws_broker.cell(row=i, column=7, value=float(getattr(b, 'target_change_pct', 0) or 0))
                ws_broker.cell(row=i, column=8, value=getattr(b, 'url', ''))
            # v4.3.2: 设置列宽
            widths = [14, 45, 16, 12, 10, 12, 12, 45]
            for i, w in enumerate(widths):
                col_letter = chr(65 + i)
                ws_broker.column_dimensions[col_letter].width = w

        # Sheet 4: 公告
        announcements = data.get("announcements", [])
        if announcements:
            ws_ann = wb.create_sheet("公告")
            headers = ["日期", "标题", "公告类型", "URL"]
            for col, h in enumerate(headers, 1):
                cell = ws_ann.cell(row=1, column=col, value=h)
                cell.font = HEADER_FONT

            for i, a in enumerate(announcements, 2):
                ws_ann.cell(row=i, column=1, value=getattr(a, 'publish_date', '')[:10])
                ws_ann.cell(row=i, column=2, value=getattr(a, 'title', ''))
                ws_ann.cell(row=i, column=3, value=getattr(a, 'announcement_type', ''))
                ws_ann.cell(row=i, column=4, value=getattr(a, 'url', ''))
            # v4.3.2: 设置列宽
            widths = [14, 50, 16, 50]
            for i, w in enumerate(widths):
                col_letter = chr(65 + i)
                ws_ann.column_dimensions[col_letter].width = w

        # v4.3.2: 批量设置所有 sheet 的数据区字体
        for ws in wb.worksheets:
            self._apply_font_to_cells(ws, start_row=2,
                                      end_row=ws.max_row,
                                      end_col=ws.max_column)

        # 保存
        if not output_path:
            output_path = str(EXPORT_DIR / f"{stock_code}_report.xlsx")

        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        wb.save(output_path)

        return output_path


class PDFExporter:
    """PDF导出器（通过Word转换）"""

    def __init__(self):
        self.word_exporter = WordExporter()

    def export_comprehensive_report(self, data: Dict[str, Any],
                                     stock_code: str,
                                     output_path: str = "") -> str:
        """
        导出综合报告为PDF
        实际通过Word文档转换实现
        """
        if not output_path:
            output_path = str(EXPORT_DIR / f"{stock_code}_report.pdf")

        # 先生成Word
        docx_path = output_path.replace('.pdf', '.docx')
        self.word_exporter.export_comprehensive_report(data, stock_code, docx_path)

        # 如果有convert工具，可以转换
        # 这里简化处理，返回docx路径并提示手动转换
        print(f"[提示] PDF导出需要手动转换: {docx_path} -> {output_path}")
        print("[提示] 可使用 LibreOffice / Word / WPS 进行转换")

        return docx_path  # 返回docx路径作为替代


class ComprehensiveExporter:
    """综合导出器 - 统一调度各格式导出"""

    def __init__(self):
        self.ppt_exporter = PPTExporter()
        self.word_exporter = WordExporter()
        self.excel_exporter = ExcelExporter()
        self.pdf_exporter = PDFExporter()
        self.organizer = DataOrganizer()
        self.analyzer = AnalysisEngine()

    def export_all_formats(self, data: Dict[str, Any],
                           stock_code: str,
                           output_dir: str = "") -> Dict[str, str]:
        """
        导出所有格式

        Args:
            data: 综合报告数据
            stock_code: 股票代码
            output_dir: 输出目录

        Returns:
            导出文件路径字典
        """
        results = {}

        if not output_dir:
            output_dir = str(EXPORT_DIR)

        Path(output_dir).mkdir(parents=True, exist_ok=True)

        stock_name = data.get("summary", {}).get("stock_name", stock_code)
        prefix = f"{output_dir}/{stock_code}_{stock_name.replace('*', '').replace('/', '')}"

        # PPT
        ppt_path = self.ppt_exporter.export_comprehensive_report(data, stock_code)
        if ppt_path:
            results["ppt"] = ppt_path

        # Word
        docx_path = self.word_exporter.export_comprehensive_report(data, stock_code)
        if docx_path:
            results["docx"] = docx_path

        # Excel
        xlsx_path = self.excel_exporter.export_comprehensive_report(data, stock_code)
        if xlsx_path:
            results["xlsx"] = xlsx_path

        # PDF (通过Word转换)
        pdf_path = self.pdf_exporter.export_comprehensive_report(data, stock_code)
        if pdf_path:
            results["pdf"] = pdf_path

        return results

    def generate_analysis_report(self, data: Dict[str, Any]) -> Dict[str, Any]:
        """
        生成分析报告

        Args:
            data: 综合报告数据

        Returns:
            分析结果
        """
        analysis = {
            "broker_rating_trend": {},
            "announcement_pattern": {},
            "summary": {}
        }

        # 分析券商研报
        broker = data.get("broker_reports", [])
        if broker:
            analysis["broker_rating_trend"] = self.analyzer.analyze_broker_rating_trend(broker)

        # 分析公告
        announcements = data.get("announcements", [])
        if announcements:
            analysis["announcement_pattern"] = self.analyzer.analyze_announcement_pattern(announcements)

        # 生成摘要
        analysis["summary"] = self.organizer.generate_summary(data)

        return analysis


# ============ CLI入口 ============

if __name__ == "__main__":
    import sys

    if len(sys.argv) < 3:
        print("用法:")
        print("  python report_exporter.py export <股票代码> --formats ppt,docx,xlsx")
        print("  python report_exporter.py analyze <股票代码>")
        print("  python report_exporter.py batch <代码1,代码2> --formats xlsx")
        sys.exit(1)

    from comprehensive_report_scraper import ComprehensiveReportManager

    cmd = sys.argv[1]

    if cmd == "export":
        code = sys.argv[2]
        formats = ["ppt", "docx", "xlsx", "pdf"]

        # 解析格式参数
        for i, arg in enumerate(sys.argv):
            if arg == "--formats" and i + 1 < len(sys.argv):
                formats = sys.argv[i + 1].split(',')

        print(f"正在获取 {code} 的数据...")
        manager = ComprehensiveReportManager()
        data = manager.get_all_reports(code)

        print(f"正在导出 {len(formats)} 种格式...")
        exporter = ComprehensiveExporter()

        for fmt in formats:
            if fmt == "ppt":
                path = exporter.ppt_exporter.export_comprehensive_report(data, code)
                print(f"  PPT: {path}")
            elif fmt == "docx":
                path = exporter.word_exporter.export_comprehensive_report(data, code)
                print(f"  Word: {path}")
            elif fmt == "xlsx":
                path = exporter.excel_exporter.export_comprehensive_report(data, code)
                print(f"  Excel: {path}")
            elif fmt == "pdf":
                path = exporter.pdf_exporter.export_comprehensive_report(data, code)
                print(f"  PDF: {path}")

        print("导出完成!")

    elif cmd == "analyze":
        code = sys.argv[2]
        manager = ComprehensiveReportManager()
        data = manager.get_all_reports(code)

        exporter = ComprehensiveExporter()
        analysis = exporter.generate_analysis_report(data)

        print(f"\n【{code} 分析报告】")
        print(f"生成时间: {analysis['summary'].get('generated_at', '')}")

        broker_trend = analysis.get("broker_rating_trend", {})
        if broker_trend:
            print(f"\n券商评级趋势:")
            print(f"  平均评级得分: {broker_trend.get('average_rating', 0):.2f}")
            print(f"  最新评级: {broker_trend.get('latest_rating', 'N/A')}")
            rating_dist = broker_trend.get("rating_distribution", {})
            if rating_dist:
                print(f"  评级分布: {rating_dist}")

        ann_pattern = analysis.get("announcement_pattern", {})
        if ann_pattern:
            print(f"\n公告模式:")
            print(f"  总数: {ann_pattern.get('total_count', 0)}")
            high_freq = ann_pattern.get("high_frequency_types", [])
            if high_freq:
                print(f"  高频类型: {', '.join(f'{k}({v})' for k, v in high_freq)}")

    elif cmd == "batch":
        if len(sys.argv) < 3:
            print("请提供股票代码（逗号分隔）")
            sys.exit(1)

        codes = sys.argv[2].split(',')
        formats = ["xlsx"]
        for i, arg in enumerate(sys.argv):
            if arg == "--formats" and i + 1 < len(sys.argv):
                formats = sys.argv[i + 1].split(',')

        manager = ComprehensiveReportManager()
        exporter = ComprehensiveExporter()

        print(f"批量导出 {len(codes)} 只股票...")

        for code in codes:
            print(f"\n处理 {code}...")
            try:
                data = manager.get_all_reports(code)
                if "xlsx" in formats:
                    path = exporter.excel_exporter.export_comprehensive_report(data, code)
                    print(f"  Excel: {path}")
            except Exception as e:
                print(f"  处理失败: {e}")

        print("\n批量导出完成!")

    else:
        print(f"未知命令: {cmd}")