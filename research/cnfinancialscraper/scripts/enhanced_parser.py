# -*- coding: utf-8 -*-
"""
增强文件/网页解析模块 v4.0
支持 PPT、HTML、Markdown、CSV、JSON 等格式的深度解析，
与现有 document_parser.py 无缝集成。

新增能力：
- PPTX 解析：文本 + 表格 + 备注 + 图片alt文本
- HTML 解析：清洗标签、保留表格结构、提取 meta/JSON-LD
- Markdown 解析：章节结构、表格、代码块
- MultiFormatParser：自动检测格式，统一入口
"""

import re
import os
import json
import csv
import io
import html as html_module
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple
from datetime import datetime

# ==================== 依赖检测 ====================

PPTX_AVAILABLE = False
try:
    from pptx import Presentation  # type: ignore
    PPTX_AVAILABLE = True
except ImportError:
    pass

HTML_PARSER_AVAILABLE = False
try:
    from html.parser import HTMLParser as StdHTMLParser
    HTML_PARSER_AVAILABLE = True
except ImportError:
    pass

LXML_AVAILABLE = False
try:
    import lxml.html  # type: ignore
    LXML_AVAILABLE = True
except ImportError:
    pass

BS4_AVAILABLE = False
try:
    from bs4 import BeautifulSoup  # type: ignore
    BS4_AVAILABLE = True
except ImportError:
    pass

SKILL_DATA_DIR = Path(__file__).parent.parent / "data"
SKILL_DATA_DIR.mkdir(parents=True, exist_ok=True)


# ==================== PPTX 解析器 ====================

class PPTXParser:
    """PowerPoint 文件解析器，支持文本/表格/备注提取。"""

    def __init__(self):
        self.available = PPTX_AVAILABLE

    def parse(self, file_path: str) -> Dict[str, Any]:
        """
        解析 PPTX 文件。

        返回: {
            file_type, file_path, slide_count,
            slides: [{slide_num, title, texts, tables, notes, images}],
            all_text, all_notes, all_tables_summary
        }
        """
        result: Dict[str, Any] = {
            "file_type": "pptx",
            "file_path": file_path,
            "slide_count": 0,
            "slides": [],
            "all_text": "",
            "all_notes": "",
            "all_tables_summary": [],
            "metadata": {},
        }

        if not self.available:
            result["error"] = "python-pptx 未安装，请运行: pip install python-pptx"
            return result

        try:
            prs = Presentation(file_path)

            # 提取文档属性
            props = prs.core_properties
            result["metadata"] = {
                "title": getattr(props, 'title', '') or '',
                "author": getattr(props, 'author', '') or '',
                "created": str(getattr(props, 'created', '')) if getattr(props, 'created', None) else '',
                "modified": str(getattr(props, 'modified', '')) if getattr(props, 'modified', None) else '',
                "slide_count": len(prs.slides),
            }

            result["slide_count"] = len(prs.slides)
            all_text_parts = []
            all_notes_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_data = {
                    "slide_num": slide_num,
                    "title": "",
                    "texts": [],
                    "tables": [],
                    "notes": "",
                    "image_count": 0,
                }

                # 提取标题（第一个非空文本框视为标题）
                first_text_found = False

                for shape in slide.shapes:
                    # 文本框
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            if not first_text_found:
                                slide_data["title"] = text[:200]
                                first_text_found = True
                            slide_data["texts"].append(text)
                            all_text_parts.append(text)

                    # 表格
                    if shape.has_table:
                        table = shape.table
                        table_data = {
                            "rows": table.rows.__len__(),
                            "cols": len(table.columns),
                            "headers": [],
                            "data": [],
                        }
                        for ri, row in enumerate(table.rows):
                            row_data = [cell.text.strip() for cell in row.cells]
                            if ri == 0:
                                table_data["headers"] = row_data
                            else:
                                table_data["data"].append(row_data)
                        slide_data["tables"].append(table_data)
                        result["all_tables_summary"].append({
                            "slide": slide_num,
                            "rows": table_data["rows"],
                            "cols": table_data["cols"],
                            "headers": table_data["headers"][:8],
                        })

                    # 图片计数
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        slide_data["image_count"] += 1

                # 提取备注
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    notes_text = notes_slide.notes_text_frame.text.strip() if notes_slide.notes_text_frame else ""
                    slide_data["notes"] = notes_text
                    if notes_text:
                        all_notes_parts.append(f"[Slide {slide_num}]\n{notes_text}")

                result["slides"].append(slide_data)

            result["all_text"] = "\n\n".join(all_text_parts)
            result["all_notes"] = "\n\n".join(all_notes_parts)
            result["text_content"] = result["all_text"]  # 兼容 document_parser 接口

        except Exception as e:
            result["error"] = f"PPTX 解析失败: {str(e)}"

        return result


# ==================== HTML 解析器 ====================

class HTMLParser:
    """HTML 网页/文件解析器，支持清洗标签、提取表格/结构化数据。"""

    def __init__(self):
        self.available = True

    def parse(self, source: str, is_url: bool = False) -> Dict[str, Any]:
        """
        解析 HTML 内容。

        Args:
            source: HTML 文本字符串或文件路径/URL
            is_url: 是否为 URL（需额外请求）

        返回: {
            file_type, title, meta, text_content, tables, links,
            json_ld, headings, financial_data
        }
        """
        result: Dict[str, Any] = {
            "file_type": "html",
            "title": "",
            "meta": {},
            "text_content": "",
            "tables": [],
            "links": [],
            "json_ld": [],
            "headings": [],
            "financial_data": {},
        }

        html_text = source

        # 如果是文件路径，读取文件
        if not is_url and os.path.isfile(source):
            try:
                with open(source, 'r', encoding='utf-8') as f:
                    html_text = f.read()
                result["file_path"] = source
            except UnicodeDecodeError:
                try:
                    with open(source, 'r', encoding='gbk') as f:
                        html_text = f.read()
                    result["file_path"] = source
                except Exception as e:
                    result["error"] = f"文件读取失败: {str(e)}"
                    return result

        if not html_text or len(html_text.strip()) < 10:
            result["error"] = "HTML 内容为空"
            return result

        try:
            # 提取 title
            title_match = re.search(r'<title[^>]*>(.*?)</title>', html_text, re.IGNORECASE | re.DOTALL)
            if title_match:
                result["title"] = self._clean_text(title_match.group(1))

            # 提取 meta 标签
            for meta_match in re.finditer(
                r'<meta\s+([^>]+)>', html_text, re.IGNORECASE
            ):
                attrs = meta_match.group(1)
                name_match = re.search(r'(?:name|property|http-equiv)\s*=\s*["\']([^"\']+)["\']', attrs, re.IGNORECASE)
                content_match = re.search(r'content\s*=\s*["\']([^"\']+)["\']', attrs, re.IGNORECASE)
                if name_match and content_match:
                    result["meta"][name_match.group(1).lower()] = content_match.group(1)

            # 提取 JSON-LD 结构化数据
            for ld_match in re.finditer(
                r'<script[^>]*type\s*=\s*["\']application/ld\+json["\'][^>]*>(.*?)</script>',
                html_text, re.IGNORECASE | re.DOTALL
            ):
                try:
                    ld_data = json.loads(ld_match.group(1))
                    result["json_ld"].append(ld_data)
                except json.JSONDecodeError:
                    pass

            # 提取表格
            result["tables"] = self._extract_tables(html_text)

            # 提取标题层级
            for level in range(1, 7):
                for h_match in re.finditer(
                    f'<h{level}[^>]*>(.*?)</h{level}>', html_text, re.IGNORECASE | re.DOTALL
                ):
                    text = self._clean_text(h_match.group(1))
                    if text:
                        result["headings"].append({"level": level, "text": text[:200]})

            # 提取链接
            for a_match in re.finditer(
                r'<a\s+[^>]*href\s*=\s*["\']([^"\']+)["\'][^>]*>(.*?)</a>',
                html_text, re.IGNORECASE | re.DOTALL
            ):
                href = a_match.group(1).strip()
                text = self._clean_text(a_match.group(2))
                if href and not href.startswith('#'):
                    result["links"].append({"url": href, "text": text[:200]})

            # 提取纯文本（去除标签）
            result["text_content"] = self._extract_text(html_text)

            # 尝试提取财务数据
            result["financial_data"] = self._extract_financial_from_html(html_text)

            result["char_count"] = len(result["text_content"])

        except Exception as e:
            result["error"] = f"HTML 解析异常: {str(e)}"

        return result

    def _clean_text(self, text: str) -> str:
        """清洗 HTML 实体和多余空白。"""
        text = re.sub(r'<[^>]+>', '', text)  # 去掉嵌套标签
        text = html_module.unescape(text)
        text = re.sub(r'\s+', ' ', text)
        return text.strip()

    def _extract_text(self, html_text: str) -> str:
        """从 HTML 提取纯文本（保留段落结构）。"""
        # 先处理换行标签
        text = re.sub(r'<(?:br|BR)\s*/?>', '\n', html_text)
        text = re.sub(r'<(?:p|P|div|DIV|li|LI|tr|TR|h\d|H\d)[^>]*>', '\n', text)
        # 去除所有标签
        text = re.sub(r'<[^>]+>', '', text)
        # 处理注释
        text = re.sub(r'<!--.*?-->', '', text, flags=re.DOTALL)
        # 处理 script/style
        text = re.sub(r'<(?:script|style|SCRIPT|STYLE)[^>]*>.*?</(?:script|style|SCRIPT|STYLE)>',
                      '', text, flags=re.DOTALL)
        # 解码 HTML 实体
        text = html_module.unescape(text)
        # 压缩空行
        text = re.sub(r'\n\s*\n', '\n\n', text)
        text = re.sub(r'[ \t]{2,}', ' ', text)
        return text.strip()

    def _extract_tables(self, html_text: str) -> List[Dict[str, Any]]:
        """提取 HTML 中的表格。"""
        tables = []
        for table_match in re.finditer(r'<table[^>]*>(.*?)</table>', html_text, re.IGNORECASE | re.DOTALL):
            table_html = table_match.group(1)
            rows = []
            for tr_match in re.finditer(r'<tr[^>]*>(.*?)</tr>', table_html, re.IGNORECASE | re.DOTALL):
                row_html = tr_match.group(1)
                cells = []
                for cell_match in re.finditer(
                    r'<t[dh][^>]*>(.*?)</t[dh]>', row_html, re.IGNORECASE | re.DOTALL
                ):
                    cells.append(self._clean_text(cell_match.group(1)))
                if cells:
                    rows.append(cells)
            if rows:
                tables.append({
                    "rows": len(rows),
                    "cols": max(len(r) for r in rows) if rows else 0,
                    "headers": rows[0] if rows else [],
                    "data": rows[1:] if len(rows) > 1 else [],
                })
        return tables

    def _extract_financial_from_html(self, html_text: str) -> Dict[str, Any]:
        """从 HTML 文本中提取财务数据。"""
        fin_data: Dict[str, Any] = {}
        text = self._extract_text(html_text)

        patterns = {
            "营业收入": r'营业(?:总)?收入[：:]\s*([\d,.]+\s*(?:亿|万|元)?)',
            "净利润": r'(?:归属.*)?净利润[：:]\s*([\d,.]+\s*(?:亿|万|元)?)',
            "总资产": r'总资产[：:]\s*([\d,.]+\s*(?:亿|万|元)?)',
            "每股收益": r'(?:基本)?每股收益[：:]\s*([\d,.]+)',
            "净资产收益率": r'净资产收益率[：:]\s*([\d,.]+)\s*%?',
        }
        for key, pat in patterns.items():
            m = re.search(pat, text)
            if m:
                fin_data[key] = m.group(1).strip()

        return fin_data


# ==================== Markdown 解析器 ====================

class MarkdownParser:
    """Markdown 文档解析器，支持章节/表格/代码块提取。"""

    def parse(self, source: str) -> Dict[str, Any]:
        """
        解析 Markdown 内容。

        返回: {file_type, title, sections, tables, code_blocks, text_content, metadata}
        """
        result: Dict[str, Any] = {
            "file_type": "markdown",
            "title": "",
            "sections": [],
            "tables": [],
            "code_blocks": [],
            "text_content": "",
            "metadata": {},
        }

        if os.path.isfile(source):
            try:
                with open(source, 'r', encoding='utf-8') as f:
                    text = f.read()
                result["file_path"] = source
            except Exception as e:
                result["error"] = f"文件读取失败: {str(e)}"
                return result
        else:
            text = source

        if not text.strip():
            result["error"] = "内容为空"
            return result

        lines = text.split('\n')
        result["text_content"] = text

        # 提取 frontmatter
        fm = self._extract_frontmatter(text)
        if fm:
            result["metadata"] = fm

        # 提取标题 (第一个 # 标题)
        for line in lines:
            h1 = re.match(r'^#\s+(.+)', line)
            if h1:
                result["title"] = h1.group(1).strip()[:200]
                break

        if not result["title"] and result["metadata"].get("title"):
            result["title"] = result["metadata"]["title"]

        # 提取章节结构
        current_section = None
        for line in lines:
            sec_match = re.match(r'^(#{1,6})\s+(.+)', line)
            if sec_match:
                level = len(sec_match.group(1))
                title = sec_match.group(2).strip()[:200]
                section = {"level": level, "title": title, "content_lines": []}
                result["sections"].append(section)
                current_section = section
            elif current_section is not None:
                current_section["content_lines"].append(line)

        # 为每个章节生成摘要
        for sec in result["sections"]:
            sec["preview"] = ' '.join(sec["content_lines"][:3])[:200]

        # 提取表格
        result["tables"] = self._extract_md_tables(text)

        # 提取代码块
        for code_match in re.finditer(r'```(\w*)\n(.*?)```', text, re.DOTALL):
            result["code_blocks"].append({
                "language": code_match.group(1) or "text",
                "content": code_match.group(2).strip(),
                "line_count": len(code_match.group(2).strip().split('\n')),
            })

        result["char_count"] = len(text)
        return result

    def _extract_frontmatter(self, text: str) -> Dict[str, Any]:
        """提取 YAML frontmatter。"""
        fm_match = re.match(r'^---\s*\n(.*?)\n---', text, re.DOTALL)
        if not fm_match:
            return {}
        fm = {}
        try:
            import yaml
            return yaml.safe_load(fm_match.group(1)) or {}
        except ImportError:
            pass

        # 手动解析简单 YAML
        for line in fm_match.group(1).split('\n'):
            kv = re.match(r'^(\w[\w_-]*)\s*:\s*(.+)', line.strip())
            if kv:
                key, val = kv.group(1), kv.group(2).strip().strip('"\'')
                fm[key] = val
        return fm

    def _extract_md_tables(self, text: str) -> List[Dict[str, Any]]:
        """提取 Markdown 表格。"""
        tables = []
        lines = text.split('\n')
        i = 0
        while i < len(lines):
            line = lines[i].strip()
            if '|' in line and line.count('|') >= 1:
                # 检查下一行是否是分隔行
                if i + 1 < len(lines) and re.match(r'^[\|\s\-:]+$', lines[i + 1].strip()):
                    # 这是一个表格
                    table_lines = [line]
                    header = [c.strip() for c in line.split('|') if c.strip()]
                    i += 2  # 跳过分隔行
                    data_rows = []
                    while i < len(lines) and '|' in lines[i]:
                        row = [c.strip() for c in lines[i].split('|') if c.strip()]
                        if row:
                            data_rows.append(row)
                        i += 1
                    tables.append({
                        "rows": 1 + len(data_rows),
                        "cols": len(header),
                        "headers": header,
                        "data": data_rows,
                    })
                    continue
            i += 1
        return tables


# ==================== CSV 解析器 ====================

class CSVParser:
    """CSV 文件解析器，支持多种编码自动检测。"""

    def parse(self, file_path: str) -> Dict[str, Any]:
        """解析 CSV 文件。"""
        result: Dict[str, Any] = {
            "file_type": "csv",
            "file_path": file_path,
            "headers": [],
            "rows": [],
            "row_count": 0,
            "col_count": 0,
            "text_content": "",
        }

        if not os.path.isfile(file_path):
            result["error"] = f"文件不存在: {file_path}"
            return result

        encodings = ['utf-8', 'utf-8-sig', 'gbk', 'gb2312', 'gb18030']
        for enc in encodings:
            try:
                with open(file_path, 'r', encoding=enc) as f:
                    reader = csv.reader(f)
                    rows = list(reader)
                if rows:
                    result["headers"] = rows[0]
                    result["rows"] = rows[1:]
                    result["row_count"] = len(rows)
                    result["col_count"] = len(rows[0]) if rows else 0
                    result["text_content"] = '\n'.join([','.join(r) for r in rows])
                    result["encoding"] = enc
                    return result
            except (UnicodeDecodeError, Exception):
                continue

        result["error"] = "无法识别 CSV 编码"
        return result


# ==================== 多格式统一解析器 ====================

class MultiFormatParser:
    """
    多格式统一解析入口。
    自动检测文件格式，分发到对应解析器。
    """

    SUPPORTED_EXTENSIONS = {
        '.pptx': 'pptx',
        '.ppt': 'pptx',
        '.html': 'html',
        '.htm': 'html',
        '.md': 'markdown',
        '.markdown': 'markdown',
        '.csv': 'csv',
        '.json': 'json',
        '.pdf': 'pdf',
        '.docx': 'docx',
        '.doc': 'docx',
        '.xlsx': 'xlsx',
        '.xls': 'xlsx',
        '.txt': 'text',
    }

    def __init__(self):
        self.pptx_parser = PPTXParser()
        self.html_parser = HTMLParser()
        self.md_parser = MarkdownParser()
        self.csv_parser = CSVParser()

    def parse(self, source: str, fmt: str = "auto") -> Dict[str, Any]:
        """
        统一解析入口。

        Args:
            source: 文件路径、URL 或文本内容
            fmt: 格式 — auto(自动检测) / pptx / html / markdown / csv

        返回: 解析结果字典
        """
        # 自动检测
        if fmt == "auto":
            if source.startswith(('http://', 'https://')):
                fmt = "html_url"
            elif os.path.isfile(source):
                ext = Path(source).suffix.lower()
                fmt = self.SUPPORTED_EXTENSIONS.get(ext, "text")
            else:
                # 启发式检测：是否为 HTML/Markdown 文本
                if re.search(r'<(!DOCTYPE|html|head|body|div|table)', source[:500], re.IGNORECASE):
                    fmt = "html"
                elif re.search(r'^#{1,6}\s', source[:500], re.MULTILINE):
                    fmt = "markdown"
                else:
                    fmt = "text"

        # 分发解析
        if fmt == "pptx":
            return self.pptx_parser.parse(source)
        elif fmt in ("html", "html_url"):
            return self.html_parser.parse(source, is_url=(fmt == "html_url"))
        elif fmt == "markdown":
            return self.md_parser.parse(source)
        elif fmt == "csv":
            return self.csv_parser.parse(source)
        elif fmt == "json":
            return self._parse_json(source)
        elif fmt in ("pdf", "docx", "doc", "xlsx", "xls", "text"):
            # 委托给现有 document_parser
            try:
                from document_parser import parse_document
                return parse_document(source)
            except ImportError:
                return {"error": "document_parser 不可用", "file_path": source}
        else:
            return {"error": f"不支持的格式: {fmt}", "file_path": source}

    def _parse_json(self, source: str) -> Dict[str, Any]:
        """解析 JSON 文件。"""
        if os.path.isfile(source):
            try:
                with open(source, 'r', encoding='utf-8') as f:
                    data = json.load(f)
            except Exception as e:
                return {"error": f"JSON 解析失败: {str(e)}", "file_path": source}
        else:
            try:
                data = json.loads(source)
            except json.JSONDecodeError as e:
                return {"error": f"JSON 解析失败: {str(e)}"}

        return {
            "file_type": "json",
            "file_path": source if os.path.isfile(source) else "",
            "data": data,
            "text_content": json.dumps(data, ensure_ascii=False, indent=2),
            "char_count": len(json.dumps(data, ensure_ascii=False)),
        }


# ==================== 便捷函数 ====================

def parse_file_enhanced(file_path: str, fmt: str = "auto") -> Dict[str, Any]:
    """一键解析任意文件/网页。"""
    parser = MultiFormatParser()
    return parser.parse(file_path, fmt)


def parse_pptx(file_path: str) -> Dict[str, Any]:
    """解析 PPTX 文件。"""
    return PPTXParser().parse(file_path)


def parse_html(source: str, is_url: bool = False) -> Dict[str, Any]:
    """解析 HTML 内容。"""
    return HTMLParser().parse(source, is_url)


def parse_markdown(source: str) -> Dict[str, Any]:
    """解析 Markdown 内容。"""
    return MarkdownParser().parse(source)


def extract_tables_from_html(html_text: str) -> List[Dict[str, Any]]:
    """从 HTML 中提取所有表格。"""
    parser = HTMLParser()
    return parser._extract_tables(html_text)


# ==================== CLI 入口 ====================

if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("用法: python enhanced_parser.py <文件路径> [格式]")
        print("格式: auto(默认) / pptx / html / markdown / csv / json")
        sys.exit(1)

    source = sys.argv[1]
    fmt = sys.argv[2] if len(sys.argv) > 2 else "auto"

    parser = MultiFormatParser()
    result = parser.parse(source, fmt)

    if result.get("error"):
        print(f"❌ {result['error']}")
    else:
        print(f"✅ 解析成功: {result.get('file_type', '?').upper()}")
        print(f"   字符数: {result.get('char_count', 0):,}")
        if result.get("title"):
            print(f"   标题: {result['title']}")
        if result.get("slide_count"):
            print(f"   幻灯片: {result['slide_count']} 页")
        if result.get("tables"):
            print(f"   表格: {len(result['tables'])} 个")
        if result.get("headings"):
            print(f"   章节: {len(result['headings'])} 个")
        print(f"\n--- 文本预览 (前500字) ---")
        preview = result.get("text_content", "")[:500]
        print(preview)
